import os
import language_tool_python
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import pandas as pd
from docx import Document
from pptx import Presentation

# Initialize grammar tool once. Check and correct grammar and spelling mistakes.
tool = language_tool_python.LanguageTool('en-US')

def show_menu():
    print("\nWhat modification would you like to apply?")
    print("1. Correct Grammar and Spelling")
    print("2. Rephrase Sentences (Simple simulation)")
    print("3. Change to Upper Case")
    print("4. Change to Lower Case")
    print("5. Change to Sentence Case")
    print("6. Capitalize Each Word")
    print("7. Toggle Case (swap upper/lower)")
    print("8. Page Formatting (add page breaks every 30 lines)")
    print("9. Create Columns (simulate with tab spaces)")
    print("10. Create Tables (basic markdown table)")
    print("11. Convert to PDF")
    print("12. Convert to Excel")
    print("13. Convert to Word")
    print("14. Convert to PowerPoint")
    print("15. Convert to Code Format (wrap in ``` )")

def modify_content(content, choice, destination_file=None):
    if choice == "1":
        matches = tool.check(content)
        return language_tool_python.utils.correct(content, matches)
    
    elif choice == "2":
        return f"**Rephrased version** (simple):\n\n{content.replace('.', '. In other words,')}"
    
    elif choice == "3":
        return content.upper()
    
    elif choice == "4":
        return content.lower()
    
    elif choice == "5":
        return '. '.join(sentence.strip().capitalize() for sentence in content.split('.'))
    
    elif choice == "6":
        return content.title()
    
    elif choice == "7":
        return content.swapcase()
    
    elif choice == "8":
        lines = content.splitlines()
        formatted = ''
        for i, line in enumerate(lines, 1):
            formatted += line + '\n'
            if i % 30 == 0:
                formatted += '\f'  # Form feed character
        return formatted
    
    elif choice == "9":
        lines = content.splitlines()
        return '\n'.join('\t'.join(line.split()) for line in lines)
    
    elif choice == "10":
        lines = content.splitlines()
        table = "| Column1 | Column2 | Column3 |\n|---------|---------|---------|\n"
        for line in lines:
            cols = line.split()
            table += f"| {' | '.join(cols[:3])} |\n"
        return table
    
    elif choice == "11":
        # Create a PDF. Convert the text into a real PDF file with page formatting.
        c = canvas.Canvas(destination_file, pagesize=letter)
        width, height = letter
        y = height - 40
        for line in content.splitlines():
            c.drawString(40, y, line)
            y -= 15
            if y < 40:
                c.showPage()
                y = height - 40
        c.save()
        return None  # No need to write text manually anymore
    
    elif choice == "12":
        # Create an Excel. Creates an Excel spreadsheet with the file content.
        data = {'Content': content.splitlines()}
        df = pd.DataFrame(data)
        df.to_excel(destination_file, index=False)
        return None
    
    elif choice == "13":
        # Create a Word document. Create a Word (.docx) document.
        doc = Document()
        for line in content.splitlines():
            doc.add_paragraph(line)
        doc.save(destination_file)
        return None
    
    elif choice == "14":
        # Create a PowerPoint. Create a basic PowerPoint presentation from provided file.
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]
        title.text = "Generated Presentation"
        body.text = '\n'.join(content.splitlines()[:10])  # Add first 10 lines
        prs.save(destination_file)
        return None
    
    elif choice == "15":
        return f"```\n{content}\n```"
    
    else:
        return content

def read_and_write_file():
    source_file = input("Enter the source filename (with extension): ")

    try:
        with open(source_file, 'r') as f:
            content = f.read()

        show_menu()
        choice = input("\nEnter the number of your choice: ")

        if choice in ["11", "12", "13", "14"]:
            destination_file = input("Enter the destination filename with proper extension (.pdf, .xlsx, .docx, .pptx): ")
        else:
            destination_file = input("Enter the destination filename to save (with extension): ")

        modified_content = modify_content(content, choice, destination_file)

        if modified_content is not None:
            with open(destination_file, 'w', encoding='utf-8') as f:
                f.write(modified_content)

        print(f"\n✅ Successfully modified and saved to '{destination_file}'.")


        #If any error occurs (missing file, bad input, etc.), the software tells the user of the error.
    except FileNotFoundError:
        print(f"❌ Error: The file '{source_file}' was not found. Please check the filename and try again.")
    except IOError:
        print(f"❌ Error: Unable to read or write files. Check file permissions or disk space.")
    except Exception as e:
        print(f"❌ An unexpected error occurred: {str(e)}")

# Run
read_and_write_file()
